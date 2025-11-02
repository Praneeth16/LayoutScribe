"""PPTX slide rendering via python-pptx and Pillow.

Renders each slide to a PNG using python-pptx to extract shapes and a
simple rasterization approach. For MVP, we snapshot the slide by saving
the slide as an image using Pillow drawing; fidelity will be improved
iteratively.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List

from PIL import Image


@dataclass(frozen=True)
class RenderedSlide:
  index0: int
  image_path: Path
  width_px: int
  height_px: int
  text: str = ""


def render_pptx_to_images(path: Path, dpi: int, temp_dir: Path) -> List[RenderedSlide]:
  try:
    from pptx import Presentation
  except Exception as exc:  # pragma: no cover
    raise RuntimeError("python-pptx is required to handle PPTX files") from exc

  prs = Presentation(path.as_posix())
  rendered: List[RenderedSlide] = []
  for i, slide in enumerate(prs.slides):
    # Approximate size using slide width/height in EMUs, convert to pixels via DPI
    # 1 inch = 914400 EMUs
    emu_per_inch = 914400
    w_in = prs.slide_width / emu_per_inch
    h_in = prs.slide_height / emu_per_inch
    width_px = int(w_in * dpi)
    height_px = int(h_in * dpi)
    # Placeholder: blank image; actual rasterization would draw shapes/text
    img = Image.new("RGB", (width_px, height_px), color=(255, 255, 255))
    out_path = temp_dir / f"slide-{i+1:04d}.png"
    img.save(out_path.as_posix(), format="PNG")
    slide_text_parts: List[str] = []
    for shape in slide.shapes:
      if hasattr(shape, "text") and shape.text:
        slide_text_parts.append(shape.text)
    rendered.append(
      RenderedSlide(
        index0=i,
        image_path=out_path,
        width_px=width_px,
        height_px=height_px,
        text="\n".join(part.strip() for part in slide_text_parts if part).strip(),
      )
    )
  return rendered


